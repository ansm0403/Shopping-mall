"""
E-커머스 프로젝트 슬라이드 2장을 기존 템플릿에 추가.
전략:
  1. 슬라이드 10 (웹툰 page1) 을 XML 레벨로 복제 -> 새 슬라이드 A (E-커머스 개요)
  2. 슬라이드 11 (웹툰 page2) 를 XML 레벨로 복제 -> 새 슬라이드 B (E-커머스 최적화)
  3. 각 textbox 의 paragraph 단위로 텍스트 교체 (첫 run 의 formatting 유지)
  4. 새 슬라이드를 10, 11 위치로 이동 (기존 웹툰은 12, 13 으로 밀림)
  5. 기술스택 슬라이드 (슬라이드 6) 의 일부 문구 교체
"""
import copy
from pptx import Presentation
from pptx.util import Pt
from lxml import etree
import os

SRC = "docs/ppt/portfolio_sangmoon.pptx"
DST = "docs/ppt/portfolio_sangmoon_v2.pptx"

prs = Presentation(SRC)

# ----------------------------------------------------------------------------
# Helper: 슬라이드 복제
# ----------------------------------------------------------------------------
def duplicate_slide(prs, index):
    """슬라이드 index 를 복제해서 끝에 추가 -> 새 슬라이드 반환"""
    src = prs.slides[index]
    blank_layout = src.slide_layout
    new_slide = prs.slides.add_slide(blank_layout)

    # 기존 placeholder 전부 제거
    for sp in list(new_slide.shapes):
        sp.element.getparent().remove(sp.element)

    # src 슬라이드의 모든 shape element 복사
    for shp in src.shapes:
        new_el = copy.deepcopy(shp.element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    return new_slide


def move_slide(prs, old_index, new_index):
    """슬라이드 순서 재배치"""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


# ----------------------------------------------------------------------------
# Helper: paragraph 의 모든 run 텍스트를 합쳐서 new_text 로 치환
#   - 첫 번째 run 의 formatting 유지
#   - 나머지 run 들은 전부 삭제
# ----------------------------------------------------------------------------
def replace_paragraph_text(paragraph, new_text):
    runs = paragraph.runs
    if not runs:
        # run 이 없으면 add
        r = paragraph.add_run()
        r.text = new_text
        return
    # 첫 run 만 남기고 나머지 삭제
    first = runs[0]
    first.text = new_text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def set_textbox_paragraphs(shape, new_paragraphs):
    """
    new_paragraphs: list[str] — 각 paragraph 의 새 텍스트
    기존 paragraph 중 len(new_paragraphs) 개만 업데이트.
    여분 paragraph 는 그대로 둔다 (빈 paragraph 들은 간격용).
    """
    tf = shape.text_frame
    paras = tf.paragraphs
    for i, text in enumerate(new_paragraphs):
        if i < len(paras):
            replace_paragraph_text(paras[i], text)
        else:
            # 새 para 추가
            p = tf.add_paragraph()
            p.text = text


def find_textbox(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def remove_picture_highlights(slide):
    """페이지2 (최적화) 슬라이드에서 노란 하이라이트 picture 를 제거
       - 페이지 배경 picture 는 남기고, 작은 사이즈 picture 만 제거"""
    # slide 11 (webtoon page2) 은 3 개의 Picture 가 있음 (배경 + 하이라이트 2개)
    # 첫 번째 picture 는 배경이므로 유지
    pics = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
    # 가장 큰 picture (배경) 제외하고 나머지 제거
    if len(pics) > 1:
        pics_sorted = sorted(pics, key=lambda p: (p.width or 0) * (p.height or 0), reverse=True)
        for p in pics_sorted[1:]:
            p.element.getparent().remove(p.element)


# ============================================================================
# 1) 슬라이드 10, 11 복제
# ============================================================================
new1 = duplicate_slide(prs, 9)   # webtoon page1 -> ecommerce page1
new2 = duplicate_slide(prs, 10)  # webtoon page2 -> ecommerce page2

# ============================================================================
# 2) 새 슬라이드 1 (E-커머스 개요) 텍스트 교체
# ============================================================================
# TextBox 4: 섹션 제목 "4. 개발경험" — 그대로 유지 (변경 없음)
# TextBox 5: 페이지 번호 "11"
tb = find_textbox(new1, "TextBox 5")
if tb:
    set_textbox_paragraphs(tb, ["11"])

# TextBox 6: 프로젝트 타이틀
tb = find_textbox(new1, "TextBox 6")
if tb:
    set_textbox_paragraphs(tb, ["#1. NestJS 기반 멀티 셀러 오픈마켓 이커머스 백엔드"])

# TextBox 7: 나의 역할
tb = find_textbox(new1, "TextBox 7")
if tb:
    set_textbox_paragraphs(tb, ["나의 역할         백엔드 API 전체 설계 및 구현 (단독 개발)"])

# TextBox 8: 팀원의 역할
tb = find_textbox(new1, "TextBox 8")
if tb:
    set_textbox_paragraphs(tb, ["팀원의 역할     없음 (개인 프로젝트)"])

# TextBox 9: 기술스택
tb = find_textbox(new1, "TextBox 9")
if tb:
    set_textbox_paragraphs(tb, [
        "기술스택",
        "",
        "백엔드 : NestJS / TypeScript / PostgreSQL / TypeORM / Redis / JWT / Docker",
        "결제/인프라 : PortOne (아임포트) / Nx 모노레포 / Throttler / Helmet",
    ])

# TextBox 10: 배포 / Github
tb = find_textbox(new1, "TextBox 10")
if tb:
    set_textbox_paragraphs(tb, [
        "배포 : 미배포 (로컬 Docker 환경에서 실행)",
        "Github : https://github.com/ansm0403/shopping_mall",
    ])

# TextBox 11: 개발 목적
tb = find_textbox(new1, "TextBox 11")
if tb:
    set_textbox_paragraphs(tb, [
        "개발 목적",
        "",
        "현재 이커머스 시장은 단순 CRUD 를 넘어 결제, 정산, 배송 추적, 역할 기반 권한 관리 등 복잡한 비즈니스 로직이 얽혀 있습니다. 이를 직접 설계하고 구현하면서 실무 수준의 백엔드 역량을 기르기 위해 기획한 프로젝트입니다.",
    ])

# TextBox 12: 프로젝트를 통해 얻은 것
tb = find_textbox(new1, "TextBox 12")
if tb:
    set_textbox_paragraphs(tb, [
        "프로젝트를 통해 얻은 것",
        "",
        "● 구매자 / 판매자 / 관리자 세 가지 역할 (RBAC) 을 중심으로 상품 승인, 판매자 신청, 주문 · 결제 · 정산까지 이커머스의 End-to-End 흐름을 직접 설계하며, 단순 CRUD 를 넘어서는 실무 수준의 도메인 모델링 경험을 쌓을 수 있었습니다.",
        "",
        "● PortOne 결제 연동 과정에서 Webhook 과 클라이언트 검증이 동시에 발생하는 동시성 문제를 트랜잭션 격리 수준 조정과 상태 체크 순서 재설계로 해결하였고, 이 과정에서 분산 환경의 결제 안정성을 확보하는 방법을 체감할 수 있었습니다.",
        "",
        "● NestJS Interceptor 와 커스텀 데코레이터 (@Auditable()) 를 활용한 AOP 기반 감사 로그 시스템을 설계하여, 기존 비즈니스 로직을 전혀 수정하지 않으면서도 61 개 이상의 주요 액션을 자동 추적하는 구조를 구현할 수 있었습니다.",
    ])

# ============================================================================
# 3) 새 슬라이드 2 (E-커머스 최적화) 텍스트 교체
# ============================================================================
# 노란 하이라이트 picture 제거 (위치가 안 맞으므로)
remove_picture_highlights(new2)

# 페이지번호 12
tb = find_textbox(new2, "TextBox 5")
if tb:
    set_textbox_paragraphs(tb, ["12"])

# TextBox 7: 프로젝트 타이틀 (page 2 상단)
tb = find_textbox(new2, "TextBox 7")
if tb:
    set_textbox_paragraphs(tb, ["#1. NestJS 기반 멀티 셀러 오픈마켓 이커머스 백엔드"])

# TextBox 6: 전체 본문 (최적화)
tb = find_textbox(new2, "TextBox 6")
if tb:
    set_textbox_paragraphs(tb, [
        "개선점 / 트러블슈팅",
        "",
        "● 결제 동시성 문제로 인한 유령 결제 해결",
        "   - PortOne Webhook 과 주문 취소가 동시에 발생할 경우 결제 상태가 불일치하며 유령 결제가 생길 수 있는 구조적 문제가 있었다.",
        "   - verifyPayment 과정에서 PortOne API 와의 동시 호출로 인한 경합 상태도 함께 존재하였다.",
        "   - 트랜잭션 격리 수준을 조정하고, 결제 검증 시 주문 상태를 먼저 확인한 뒤 처리하는 순서로 재설계하여 유령 결제 발생 건수를 0 건으로 줄일 수 있었다.",
        "",
        "● 감사 로그 시스템 AOP 설계",
        "   - 이커머스 특성상 결제, 주문, 환불 등 민감한 동작이 많아 추적이 필요했으나, 각 서비스마다 수동 로깅 코드를 삽입하면 비즈니스 로직이 오염되는 문제가 있었다.",
        "   - NestJS Interceptor 와 @Auditable() 커스텀 데코레이터를 조합하여 기존 서비스 코드 0 줄 수정으로 61 개 이상의 액션을 자동 추적하는 구조를 설계하였다.",
        "   - IP 주소, User-Agent, 성공 / 실패 여부까지 기록하여 관리자가 감사 로그를 실시간으로 조회할 수 있도록 하였다.",
        "",
        "● 멀티 셀러 배송 분리 설계",
        "   - 하나의 주문에 여러 판매자의 상품이 포함될 수 있어, 주문 (Order) 단위로 배송 상태를 관리할 경우 판매자별 독립 처리가 불가능한 구조적 한계가 있었다.",
        "   - 배송 (Shipment) 엔티티를 주문이 아닌 판매자 (Seller) 단위로 분리하여 각 판매자가 독립적으로 배송 처리 및 운송장 등록을 할 수 있도록 개선하였다.",
    ])

# ============================================================================
# 4) 슬라이드 순서 재배치
#   현재: 0..9 = 기존 1~10번 슬라이드, 10 = 기존 슬라이드 11
#   11, 12 = 기존 slide 12, 13 (일본호텔), 13,14 = 기존 slide 14,15 (계좌), 15 = 끝
#   16, 17 = 방금 추가된 E-커머스 page1, page2
#   목표: E-커머스를 슬라이드 10, 11 위치로 이동 (웹툰 앞이 아니라 웹툰 자리)
#   즉, 기존 웹툰 slide (index 9, 10) 은 뒤로 밀려야 함
# ============================================================================
# 현재 slides index 순서: [0..15 기존, 16=new1(ecom p1), 17=new2(ecom p2)]
# 목표: ecom 을 9, 10 위치로 이동 (기존 9,10 = 웹툰이 11, 12 로 밀림)
move_slide(prs, 16, 9)   # ecom p1 -> 9
move_slide(prs, 17, 10)  # ecom p2 -> 10 (방금 9로 이동한 후 17은 여전히 끝에 있음)

# ============================================================================
# 5) 기술스택 슬라이드 (슬라이드 6, index 5) 업데이트
#   "프론트엔드" 박스를 유지하되 일부 프론트/백엔드 혼합
#   Sanity -> PostgreSQL, Cloud firestore -> Redis, Firebase -> Docker 등으로 교체
# ============================================================================
tech_slide = prs.slides[5]

# 기술스택 슬라이드의 텍스트 박스를 이름으로 찾되, 내용으로 매칭해서 교체
# 각 TextBox 내용 치환 규칙 (content-based replace)
replace_map = {
    "Sanity": "PostgreSQL",
    "Sanity를 통해 댓글과 좋아요 등의 커뮤니티 기능을 구현함으로써, 페이지사이의 기능을 추가했습니다.":
        "관계형 DB 로 이커머스의 복잡한 도메인 (주문, 결제, 정산, 감사 로그 등) 을 모델링하고, TypeORM 과 결합하여 트랜잭션 기반의 안정적인 데이터 처리를 구현하였습니다.",
    "Cloud firestore": "Redis",
    "Firebase의 Firestore를 활용하여 실시간 데이터 업데이트 기능과 제작 애플리케이션의 프런트엔드를 구현했습니다.":
        "토큰 블랙리스트 및 다중 디바이스 세션 관리, Rate Limiting 등 휘발성이 강한 상태 데이터를 빠르게 처리하기 위해 인메모리 캐시로 활용하였습니다.",
    "prisma + supabase": "Docker / AWS",
    "숙련도가 낮은편이지만 간단한 프로젝트에서 공공데이터를 이용하여 데이터 관리 및 댓글, 좋아요 기능을 구현하였습니다.":
        "Docker Compose 로 PostgreSQL · Redis 등 인프라를 컨테이너 환경으로 구성하여 개발 · 배포의 일관성을 확보하였으며, AWS 배포 경험을 쌓아가고 있습니다.",
}

for shape in tech_slide.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        full_text = "".join(r.text for r in para.runs)
        stripped = full_text.strip()
        if stripped in replace_map:
            replace_paragraph_text(para, replace_map[stripped])

# ============================================================================
# 저장
# ============================================================================
prs.save(DST)
print(f"Saved: {DST}")
print(f"Total slides: {len(prs.slides)}")
